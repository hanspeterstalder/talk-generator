import pickle
import random
import pathlib
import os.path
import inflect
import argparse
import requests
import safygiphy
import math
import numpy
from random import randint
from os import listdir
from os.path import isfile, join

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from bs4 import BeautifulSoup

import nltk
from nltk.corpus import wordnet as wn
from py_thesaurus import Thesaurus
from google_images_download import google_images_download

# CONSTANTS
HEIGHT = 9
WIDTH = 16
LEFTMOST = Inches(0)
TOPMOST = Inches(0)
HEIGHT_IN = Inches(HEIGHT)
WIDTH_IN = Inches(WIDTH)

# One inch equates to 914400 EMUs 
INCHES_TO_EMU = 914400
# One centimeter is 360000 EMUs
CMS_TO_EMU = 360000

# Location of powerpoint template
POWERPOINT_TEMPLATE_FILE = 'data/powerpoint/template.pptx'


# Slide generator class
class SlideGenerator:

    # Class function to create a function that always returns a certain weight
    def constant_weight(weight: int):
        return lambda slide_nr, total_slides: weight

    def __init__(self, generator, weight_function=constant_weight(1)):
        self._generator = generator
        self._weight_function = weight_function

    # Generate a slide for a given presentation using the given seed.
    def generate(self, presentation, seed):
        slide = self._generator(presentation, seed)
        # Add information about the generator to the notes
        if slide:
            slide.notes_slide.notes_text_frame.text = str(self)

    # The weight of the generator for a particular slide.
    # Determines how much chance it has being picked for a particular slide number
    def get_weight_for(self, slide_nr, total_slides):
        return self._weight_function(slide_nr, total_slides)

    def __str__(self):
        return "SlideGenerator[" + str(self._generator.__name__) + "]"


# Class responsible for determining which slide generators to use in a presentation, and how the (topic) seed for
# each slide is generated
class PresentationSchema:

    def __init__(self, seed_generator, slide_generators):
        self._seed_generator = seed_generator
        self._slide_generators = slide_generators

    # Generate a presentation about a certain topic with a certain number of slides
    def generate_presentation(self, topic, number_of_slides):

        # Create new presentation
        presentation = Presentation(POWERPOINT_TEMPLATE_FILE)
        # Create the topic-for-each-slide generator
        seed_generator = self._seed_generator(topic, number_of_slides)

        for i in range(number_of_slides):
            # Generate a topic for the next slide
            seed = seed_generator.generate_seed(i)

            # Select the slide generator to generate with
            generator = self._select_generator(i, number_of_slides)

            # Generate the slide
            print('Adding slide {} about {} using {}'.format(i, seed, generator))
            slide = generator.generate(presentation, seed)

            # Try again if slide is False!
            if not bool(slide):
                i -= 1
                # TODO: Remove slide from presentation if there was a slide generated

        return presentation

    # Select a generator for a certain slide number
    def _select_generator(self, slide_nr, total_slides):
        weighted_generators = []
        for i in range(len(self._slide_generators)):
            generator = self._slide_generators[i]
            weighted_generator = generator.get_weight_for(slide_nr, total_slides), generator
            weighted_generators.append(weighted_generator)
        return weighted_random(weighted_generators)


# This class generates a bunch of related words (e.g. synonyms) of a word to generate topics for a presentation
class SynonymTopicGenerator:

    def __init__(self, topic, number_of_slides):
        self._topic = topic
        self._slides_nr = number_of_slides
        seeds = get_synonyms(topic)
        # seeds.extend(get_relations(topic))

        # Check if enough generated
        if len(seeds) < number_of_slides:
            # If nothing: big problem!
            if len(seeds) == 0:
                seeds = [topic]

            # Now fill the seeds up with repeating topics
            number_of_repeats = int(math.ceil(number_of_slides / len(seeds)))
            seeds = numpy.tile(seeds, number_of_repeats)

        # Take random `number_of_slides` elements
        random.shuffle(seeds)
        self._seeds = seeds[0: number_of_slides]

    def generate_seed(self, slide_nr):
        return self._seeds[slide_nr]


# HELPER FUNCTIONS
def _save_presentation_to_pptx(args, prs):
    """Save the talk."""
    fp = './output/' + args.topic + '.pptx'
    # Create the parent folder if it doesn't exist
    pathlib.Path(os.path.dirname(fp)).mkdir(parents=True, exist_ok=True)
    prs.save(fp)
    print('Saved talk to {}'.format(fp))
    return True


def download_image(fromUrl, toUrl):
    """Download image from url to path."""
    # Create the parent folder if it doesn't exist
    pathlib.Path(os.path.dirname(toUrl)).mkdir(parents=True, exist_ok=True)

    # Download
    f = open(toUrl, 'wb')
    f.write(requests.get(fromUrl).content)
    f.close()


def read_lines(file):
    return [line.rstrip('\n') for line in open(file)]


# From https://stackoverflow.com/questions/14992521/python-weighted-random
def weighted_random(pairs):
    total = sum(pair[0] for pair in pairs)
    r = randint(1, total)
    for (weight, value) in pairs:
        r -= weight
        if r <= 0: return value


# CONTENT GENERATORS
# These functions generate content, sometimes related to given arguments

def get_definitions(word):
    """Get definitions of a given topic word."""
    print('******************************************')
    # Get definition
    word_senses = wn.synsets(word)
    definitions = {}
    for ss in word_senses:
        definitions[ss.name()] = ss.definition()
    print('{} definitions for "{}"'.format(len(definitions), word))
    return definitions


def get_synonyms(word):
    """Get all synonyms for a given word."""
    print('******************************************')
    word_senses = wn.synsets(word)
    all_synonyms = []
    for ss in word_senses:
        all_synonyms.extend(
            [x.lower().replace('_', ' ') for x in ss.lemma_names()])
    all_synonyms.append(word)
    all_synonyms = list(set(all_synonyms))
    print('{} synonyms for "{}"'.format(len(all_synonyms), word))
    return all_synonyms


def get_relations(word):
    """Get relations to given definitions."""
    rels = {}
    all_rel_forms = []
    all_perts = []
    all_ants = []

    word_senses = wn.synsets(word)
    for ss in word_senses:
        ss_name = ss.name()
        rels[ss_name] = {}
        for lem in ss.lemmas():
            lem_name = lem.name()
            rels[ss_name][lem_name] = {}
            rel_forms = [x.name() for x in lem.derivationally_related_forms()]
            rels[ss_name][lem_name]['related_forms'] = rel_forms
            all_rel_forms.extend(rel_forms)

            perts = [x.name() for x in lem.pertainyms()]
            rels[ss_name][lem_name]['pertainyms'] = perts
            all_perts.extend(perts)

            ants = [x.name() for x in lem.antonyms()]
            rels[ss_name][lem_name]['antonyms'] = ants
            all_ants.extend(ants)

    print('******************************************')
    print('{} derivationally related forms'.format(len(all_rel_forms)))
    print('******************************************')
    print('{} pertainyms'.format(len(all_perts)))
    print('******************************************')
    print('{} antonyms'.format(len(all_ants)))
    return rels


def get_title(synonyms):
    """Returns a template title from a source list."""
    print('******************************************')
    chosen_synonym = random.choice(synonyms)
    chosen_synonym_plural = inflect.engine().plural(chosen_synonym)
    synonym_templates = read_lines('data/text-templates/titles.txt')
    chosen_template = random.choice(synonym_templates);
    return chosen_template.format(chosen_synonym_plural.title())


def get_images(synonyms, num_images):
    """Get images, first search locally then Google Image Search."""
    all_paths = {}
    if num_images > 0:
        for word in synonyms:
            all_paths[word] = get_google_images(word, num_images)

    return all_paths


def get_google_images(word, num_images=1):
    lp = 'downloads/' + word + '/'
    paths = _get_google_image_cached(word, num_images, lp)

    # If no local images, search on Google Image Search
    if len(paths) == 0:
        # Get related images at 16x9 aspect ratio
        response = google_images_download.googleimagesdownload()
        arguments = {
            'keywords': word,
            'limit': num_images,
            'print_urls': True,
            'exact_size': '1600,900',
        }
        # passing the arguments to the function
        paths = response.download(arguments)
        # printing absolute paths of the downloaded images
        print('paths of images', paths)
    return paths


def _get_google_image_cached(word, num_image, lp):
    paths = []
    try:
        local_files = [lp + f for f in listdir(lp) if isfile(join(lp,
                                                                  f))]
        paths = local_files
    except FileNotFoundError as e:
        paths = []

    if len(paths) > 0:
        print('{} local images on {} found'.format(len(paths), word))

    return paths


def get_related_giphy(seed_word):
    giphy = safygiphy.Giphy()
    result = giphy.random(tag=seed_word)
    return result.get('data').get('images').get('original').get('url')


def wikihow_action_to_action(wikihow_title):
    index_of_to = wikihow_title.find('to')
    return wikihow_title[index_of_to + 3:]


def search_wikihow(search_words):
    return requests.get(
        'https://en.wikihow.com/wikiHowTo?search='
        + search_words.replace(' ', '+'))


def get_related_wikihow_actions(seed_word):
    page = search_wikihow(seed_word)
    # Try again but with plural if nothing is found
    if not page:
        page = search_wikihow(inflect.engine().plural(seed_word))

    soup = BeautifulSoup(page.content, 'html.parser')
    actions_elements = soup.find_all('a', class_='result_link')
    actions = \
        list(
            map(wikihow_action_to_action,
                map(lambda x: x.get_text(), actions_elements)))

    return actions


# FORMAT GENERATORS
# These are functions that get some inputs (texts, images...) 
# and create layouted slides with these inputs

def create_title_slide(args, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_object = slide.shapes.title
    title_object.text = args.title
    return slide


def create_text_slide(prs, text):
    # Get a default blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_object = slide.shapes.title
    title_object.text = text
    return slide


# Creates a slide with an image covering the whole slide
def create_image_slide(prs, image_url, title=False):
    return _create_single_image_slide(prs, image_url, 2, title)


# Creates a slide with an image covering the whole slide
def create_full_image_slide(prs, image_url, title=False):
    return _create_single_image_slide(prs, image_url, 11, title)


def _create_single_image_slide(prs, image_url, slide_template_idx, title=False):
    # Add image url as picture
    if image_url:
        # Get a default blank slide layout
        slide = prs.slides.add_slide(prs.slide_layouts[slide_template_idx])

        if title:
            title_object = slide.shapes.title
            title_object.text = title

        image_placeholder = slide.placeholders[1]
        image_placeholder.insert_picture(image_url)

        return slide

    return False


# FULL SLIDES GENERATORS:
# These are functions that create slides with certain (generated) content

def create_google_image_slide(prs, seed_word):
    # Get all image paths
    # img_paths = args.all_paths.get(word)
    img_paths = get_google_images(seed_word, 1)
    if img_paths:
        print("PATHS:" + str(img_paths))
        # Pick one of the images
        img_path = random.choice(img_paths)

        # Create slide with image
        slide = create_full_image_slide(prs, img_path, seed_word)

        # Add title to the slide
        if bool(slide):
            shapes = slide.shapes
            shapes.title.text = seed_word
            return slide
    return False


def create_inspirobot_slide(prs, topic):
    # Generate a random url to access inspirobot
    dd = str(random.randint(1, 73)).zfill(2)
    nnnn = random.randint(0, 9998)
    inspirobot_url = ('http://generated.inspirobot.me/'
                      '0{}/aXm{}xjU.jpg').format(dd, nnnn)

    # Download the image
    image_url = 'downloads/inspirobot/{}-{}.jpg'.format(dd, nnnn)
    download_image(inspirobot_url, image_url)

    print("Downloaded inspirobot image: {}".format(image_url))

    # Turn into image slide
    return create_full_image_slide(prs, image_url)


def create_giphy_slide(prs, word):
    # Download the image
    giphy_url = get_related_giphy(word)
    gif_name = os.path.basename(os.path.dirname(giphy_url))
    image_url = 'downloads/' + word + '/gifs/' + gif_name + ".gif"
    download_image(giphy_url, image_url)

    # Turn into image slide
    return create_full_image_slide(prs, image_url)


def create_wikihow_action_bold_statement_slide(prs, wikihow_seed):
    related_actions = get_related_wikihow_actions(wikihow_seed)
    if related_actions:
        action = random.choice(related_actions)
        bold_statement_templates = read_lines('data/text-templates/bold-statements.txt')

        chosen_template = random.choice(bold_statement_templates)
        template_values = {'action': action.title(),
                           # TODO: Make a scraper that scrapes a step related to this action on wikihow.
                           # TODO: Fix action_infinitive
                           'action_infinitive': action.title(),
                           'step': 'Do Whatever You Like',
                           'topic': wikihow_seed,
                           # TODO: Use datamuse or some other mechanism of finding a related location
                           'location': 'Here'}
        life_lesson = chosen_template.format(**template_values)

        # Turn into image slide
        return create_text_slide(prs, life_lesson)
    return False


# COMPILATION
# Compiling the slides to a powerpoint


def compile_talk_to_pptx(args):
    """Compile the talk with the given source material."""
    prs = Presentation(POWERPOINT_TEMPLATE_FILE)
    # Set the height and width
    # prs.slide_height = HEIGHT * INCHES_TO_EMU
    # prs.slide_width = WIDTH * INCHES_TO_EMU

    # Build an ordered list of slides for access
    slides = []

    # Add title slide
    title_slide = create_title_slide(args, prs)
    slides.append(title_slide)
    slide_idx_iter = 1

    # For each synonym 
    for word, path_list in args.all_paths.items():
        # print('Word: {}'.format(word))
        # For each image collected add a new slide
        # for i in range(len(path_list)):
        print('***********************************')
        print('Adding slide {} about {}'.format(slide_idx_iter, word))
        slide = create_google_image_slide(args, prs, word)
        if slide:
            slides.append(slide)
            slide_idx_iter += 1

    # Add some Inspirobot quotes
    print('***********************************')
    print('Adding slide: {}, Inspirobot'.format(slide_idx_iter))
    slide = create_inspirobot_slide(prs)
    if slide:
        slides.append(slide)
        slide_idx_iter += 1

    # Add a Gif slide
    print('***********************************')
    giphy_seed = random.choice(args.synonyms)
    print('Adding slide: {}, Giphy about {}'.format(slide_idx_iter, giphy_seed))
    slide = create_giphy_slide(prs, giphy_seed)
    if slide:
        slides.append(slide)
        slide_idx_iter += 1

    # Add a life lesson
    print('***********************************')
    for i in range(2):
        wikihow_seed = random.choice(args.synonyms)
        print('Adding Wikihow Lifelesson slide: {} about {}'.format(slide_idx_iter,
                                                                    wikihow_seed))
        slide = create_wikihow_action_bold_statement_slide(prs, wikihow_seed)
        if bool(slide):
            slides.append(slide)
            slide_idx_iter += 1

    _save_presentation_to_pptx(args, prs)
    print('Successfully built talk.')


def compile_talk_to_raw_data(args):
    """Save the raw data that has been harvested."""
    with open('output/' + args.topic.replace(' ', '_') + '.pkl', 'wb') as fh:
        pickle.dump(args, fh, protocol=pickle.HIGHEST_PROTOCOL)
        print('Pickle saved to output/' + args.topic.replace(' ', '_') + '.pkl')


# MAIN

def main(args):
    """Make a talk with the given topic."""
    # Print status details
    print('******************************************')
    print("Making {} slide talk on: {}".format(args.num_slides, args.topic))

    # Parse topic string to parts-of-speech
    text = nltk.word_tokenize(args.topic)
    print('******************************************')
    print('tokenized text: ', text)
    print('pos tag text: ', nltk.pos_tag(text))

    # Parse the actual topic subject from the parts-of-speech
    topic_string = args.topic

    # Get definitions
    args.definitions = get_definitions(topic_string)
    # Get relations
    args.relations = get_relations(topic_string)
    # Get synonyms
    args.synonyms = get_synonyms(topic_string)
    # Get related actions
    args.actions = get_related_wikihow_actions(topic_string)
    # Get a title
    args.title = get_title(args.synonyms)
    # For each synonym download num_images
    args.all_paths = get_images(args.synonyms, args.num_images)

    # Compile and save the presentation to data
    compile_talk_to_raw_data(args)

    # Compile and save the presentation to PPTX
    # compile_talk_to_pptx(args)
    presentation = presentation_schema.generate_presentation(topic_string, args.num_slides)

    # Save presentation
    _save_presentation_to_pptx(args, presentation)


presentation_schema = PresentationSchema(lambda topic, num_slides: SynonymTopicGenerator(topic, num_slides),
                                         [
                                             SlideGenerator(create_giphy_slide),
                                             SlideGenerator(create_inspirobot_slide),
                                             SlideGenerator(create_wikihow_action_bold_statement_slide),
                                             SlideGenerator(create_google_image_slide)])

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Process some integers.')
    parser.add_argument('--topic', help="Topic of presentation.",
                        default='bagels', type=str)
    parser.add_argument('--num_images', help="Number of images per synonym.",
                        default=1, type=int)
    parser.add_argument('--num_slides', help="Number of slides to create.",
                        default=3, type=int)
    args = parser.parse_args()
    main(args)
